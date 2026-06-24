import os
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from utils.exceptions import FileParseError
from utils.logger import setup_logger
from utils.glossary import glossary_file_has_entries
from utils.industry import label_for_industry_id

logger = setup_logger()


def is_target_only_layout(fmt: str) -> bool:
    """纯译文输出（兼容旧配置 chinese_only）。"""
    return fmt in ("chinese_only", "target_only")


def compute_translation_output_path(
    file_path: str,
    *,
    source_lang: str = "en",
    target_lang: str = "zh-Hans",
    output_layout: str = "bilingual",
    glossary_file_path: str = "",
    industry_preset: str = "general",
) -> str:
    """
    与 FileGenerator._generate_output_path 规则一致，供批量扫描时判断译稿是否已生成。
    file_path 应为实际参与命名的工作文件路径（与生成阶段传入生成器的路径一致）。
    """
    path = Path(file_path)
    src = _safe_filename_lang_segment(source_lang)
    tgt = _safe_filename_lang_segment(target_lang)
    ind_raw = label_for_industry_id((industry_preset or "general").strip() or "general")
    ind_token = _safe_filename_glossary_stem(ind_raw, max_len=40)
    if not ind_token:
        ind_token = "通用"
    if is_target_only_layout(output_layout):
        layout_token = "纯译文"
    else:
        layout_token = "双语对照"
    parts = [path.stem, src, tgt, ind_token, layout_token]
    gpath = (glossary_file_path or "").strip()
    if gpath and glossary_file_has_entries(gpath):
        gstem = _safe_filename_glossary_stem(Path(gpath).stem)
        if gstem:
            parts.append(gstem)
    parts.append("XuYuan")
    filename = "_".join(parts) + path.suffix
    return str(path.parent / filename)


def _safe_filename_lang_segment(code: str) -> str:
    """语言代码写入文件名时去掉非法字符。"""
    s = (code or "xx").strip()
    bad = '<>:"/\\|?*\n\r\t'
    out = "".join("_" if c in bad else c for c in s)
    out = out.strip(" .")
    return out or "lang"


def _safe_filename_glossary_stem(stem: str, max_len: int = 40) -> str:
    """词汇表文件名（不含扩展名）写入输出文件名时净化并限长。"""
    s = (stem or "").strip()
    bad = '<>:"/\\|?*\n\r\t'
    out = "".join("_" if c in bad else c for c in s)
    out = out.replace(" ", "_")
    while "__" in out:
        out = out.replace("__", "_")
    out = out.strip(" ._")
    if not out:
        return ""
    return out[:max_len].rstrip("_")


def _ppt_avg_char_width_pt(font_pt: float, text: str) -> float:
    """混合文本下单字符平均占位宽度（pt，经验值，用于估算换行）。"""
    if not text:
        return max(font_pt * 0.45, 1.0)
    wide = 0
    for c in text:
        o = ord(c)
        if 0x4E00 <= o <= 0x9FFF:
            wide += 1
        elif 0x3040 <= o <= 0x30FF or 0x31F0 <= o <= 0x31FF:
            wide += 1
        elif 0xAC00 <= o <= 0xD7A3:
            wide += 1
        elif 0x0600 <= o <= 0x06FF or 0x0750 <= o <= 0x077F:
            wide += 1
    latin = len(text) - wide
    if latin + wide == 0:
        return font_pt * 0.5
    return (wide * font_pt * 0.92 + latin * font_pt * 0.48) / len(text)


def _ppt_latin_letter_ratio(text: str) -> float:
    """拉丁字母占比（对照后英文变长时用于略收紧换行估算）。"""
    if not text:
        return 0.0
    n = len(text)
    lat = sum(1 for c in text if ("A" <= c <= "Z") or ("a" <= c <= "z"))
    return lat / n


def _ppt_wrapped_content_height_pt(
    text: str,
    font_pt: float,
    box_w_pt: float,
    line_factor: float = 1.24,
    *,
    width_scale: float = 1.0,
    height_pessimism: float = 1.0,
) -> float:
    """按固定栏宽估算块文本换行后的总高度（pt）。"""
    if font_pt <= 0 or box_w_pt <= 0:
        return 1e9
    eff_w = max(box_w_pt * width_scale, font_pt * 2)
    acw = max(font_pt * 0.35, _ppt_avg_char_width_pt(font_pt, text) * 1.06)
    chars_per_line = max(3, int(eff_w / acw))
    n = len(text)
    lines = max(1, (n + chars_per_line - 1) // chars_per_line)
    return lines * font_pt * line_factor * height_pessimism


def _ppt_text_frame_inner_box_pt(text_frame, shape) -> Tuple[Optional[float], Optional[float]]:
    """文本框内可排版区域约宽×高（pt）。"""
    try:
        tw = shape.width.pt - text_frame.margin_left.pt - text_frame.margin_right.pt
        th = shape.height.pt - text_frame.margin_top.pt - text_frame.margin_bottom.pt
        return max(10.0, tw), max(10.0, th)
    except Exception:
        return None, None


def _ppt_paragraph_target_heights_pt(text_frame, shape) -> Dict[int, float]:
    """
    同一 text_frame 内多段（项目符号）共享垂直空间时，按各段「原文」估算排版高度比例，
    为每段分配可用的目标高度（对照/仅译文均需要），避免每段都按整框高度算字号导致严重溢出。
    """
    w_pt, h_pt = _ppt_text_frame_inner_box_pt(text_frame, shape)
    if w_pt is None or h_pt is None or w_pt < 16 or h_pt < 14:
        return {}

    indexed = [
        (i, p) for i, p in enumerate(text_frame.paragraphs) if p.text and p.text.strip()
    ]
    if not indexed:
        return {}
    if len(indexed) == 1:
        idx = indexed[0][0]
        return {idx: h_pt * 0.82}

    usable = h_pt * 0.74
    weights: List[float] = []
    for _, p in indexed:
        ot = p.text.strip()
        orig_pt = 12.0
        if p.runs and p.runs[0].font.size:
            try:
                orig_pt = float(p.runs[0].font.size.pt)
            except Exception:
                pass
        wh = _ppt_wrapped_content_height_pt(
            ot, orig_pt, w_pt, line_factor=1.28
        )
        weights.append(max(wh, orig_pt * 1.08))

    tot = sum(weights)
    if tot <= 0:
        n = len(indexed)
        return {indexed[k][0]: usable / n for k in range(n)}

    return {
        indexed[k][0]: usable * (weights[k] / tot) for k in range(len(indexed))
    }


def _ppt_equal_paragraph_target_heights_pt(
    text_frame, shape
) -> Dict[int, float]:
    """
    同一文本框内多条非空段落均分可用垂直高度。
    避免「按比例按原文高度加权」时首段占满预算、后段 target_h 过小导致字号明显变小。
    """
    non_empty = [
        i
        for i, p in enumerate(text_frame.paragraphs)
        if p.text and p.text.strip()
    ]
    if len(non_empty) <= 1:
        return {}
    w_pt, h_pt = _ppt_text_frame_inner_box_pt(text_frame, shape)
    if w_pt is None or h_pt is None or w_pt < 16 or h_pt < 14:
        return {}
    # 仅用于「宽幅正文」形状；与 _ppt_paragraph_target_heights_pt 中 usable 比例接近
    usable = h_pt * 0.76
    eh = usable / float(len(non_empty))
    return {i: eh for i in non_empty}


def _ppt_frame_max_run_font_pt(text_frame) -> float:
    """取该文本框各非空段首 run 字号的最大值，作同框统一缩放时的上限参考。"""
    m = 12.0
    for p in text_frame.paragraphs:
        if not p.text or not p.text.strip():
            continue
        if p.runs and p.runs[0].font.size:
            try:
                m = max(m, float(p.runs[0].font.size.pt))
            except Exception:
                pass
    return m


def _ppt_shape_is_wide_body_text_block(text_frame, shape) -> bool:
    """
    判断是否适合应用「多段均分高度 + 统一字号 + 放松」的宽幅正文框。
    窄框（图示标签、流程图内文字等）必须为 False，否则会强行放大字号导致竖排叠字、重叠。
    """
    w_pt, h_pt = _ppt_text_frame_inner_box_pt(text_frame, shape)
    if w_pt is None or h_pt is None:
        return False
    if w_pt < 200.0:
        return False
    if h_pt < 36.0:
        return False
    return True


def _ppt_readability_floor_pt(
    orig_pt: float,
    new_text: str,
    original_text: str,
    only_translation: bool,
) -> float:
    """
    在防溢出计算之后，尽量抬升「可读下限」，避免标签、短句、小框说明被压得过小。
    若抬升后仍由后续 _ppt_wrapped_content_height_pt 校验可放入 target_h 才采用。
    """
    nt = (new_text or "").strip()
    ot = (original_text or "").strip()
    n = len(nt)
    ro = max(len(ot), 1)
    ratio_len = n / ro

    hard_min = 7.25 if orig_pt >= 9.5 else 6.5

    if n <= 24:
        return max(hard_min, min(orig_pt, orig_pt * 0.90))
    if n <= 50:
        return max(hard_min, orig_pt * 0.82)
    if n <= 95:
        return max(hard_min, orig_pt * 0.72)

    if ratio_len <= 1.32:
        return max(hard_min, orig_pt * 0.84)
    if ratio_len <= 1.65:
        return max(hard_min, orig_pt * 0.74)

    return max(hard_min, orig_pt * (0.58 if only_translation else 0.54))


def _ppt_bilingual_fit_font_pt(
    orig_pt: float,
    original_text: str,
    new_text: str,
    text_frame,
    shape,
    target_h_override: Optional[float] = None,
    *,
    only_translation: bool = False,
) -> float:
    """
    在文本框近似版面内选最大字号（不超过原字号），使内容尽量落在框内。
    对照模式与仅译文模式共用；仅译文常为「中文原文→更长英文」，略额外收紧。
    target_h_override：本段在共享文本框内分到的可用高度（pt）；未给则用整框高度。
    """
    min_pt = 6.5
    try:
        o = float(orig_pt)
    except (TypeError, ValueError):
        o = 12.0
    o = max(min_pt, min(o, 120.0))

    w_pt, h_pt = _ppt_text_frame_inner_box_pt(text_frame, shape)
    lr = _ppt_latin_letter_ratio(new_text or "")
    nt_for_len = (new_text or "").strip()
    if only_translation:
        width_scale = 0.87 if lr > 0.32 else (0.91 if lr > 0.16 else 0.94)
        height_pessimism = 1.11 if lr > 0.24 else 1.06
        line_factor = 1.30 if lr > 0.20 else 1.26
    else:
        width_scale = 0.88 if lr > 0.38 else (0.92 if lr > 0.22 else 0.96)
        height_pessimism = 1.12 if lr > 0.28 else 1.07
        line_factor = 1.32 if lr > 0.25 else 1.27

    # 窄框且正文较长时再收紧；短标签避免无谓压栏宽
    if w_pt is not None and w_pt < 72 and len(nt_for_len) > 55:
        width_scale *= 0.94

    if target_h_override is not None and target_h_override > 0:
        target_h = max(min_pt * 2.0, float(target_h_override))
        if only_translation:
            target_h *= 0.945
    else:
        target_h = (h_pt * 0.82) if h_pt else None
        if only_translation and target_h:
            target_h *= 0.965

    if (
        w_pt is not None
        and h_pt is not None
        and target_h
        and w_pt >= 16
        and h_pt >= 14
    ):
        lo, hi = min_pt, o
        best = min_pt
        for _ in range(28):
            mid = (lo + hi) / 2.0
            est = _ppt_wrapped_content_height_pt(
                new_text,
                mid,
                w_pt,
                line_factor=line_factor,
                width_scale=width_scale,
                height_pessimism=height_pessimism,
            )
            if est <= target_h:
                best = mid
                lo = mid
            else:
                hi = mid
            if hi - lo < 0.07:
                break
        floor_pt = _ppt_readability_floor_pt(
            o, new_text, original_text, only_translation
        )
        floor_pt = min(floor_pt, o)
        if floor_pt > best + 0.04:
            est_floor = _ppt_wrapped_content_height_pt(
                new_text,
                floor_pt,
                w_pt,
                line_factor=line_factor,
                width_scale=width_scale,
                height_pessimism=height_pessimism,
            )
            if est_floor <= target_h * 1.02:
                best = floor_pt
        return max(min_pt, min(best, o))

    ro = max(len((original_text or "").strip()), 1)
    rn = max(len(new_text or ""), 1)
    ratio = rn / ro
    if ratio <= 1.12:
        return o
    exp = 0.38 if only_translation else 0.40
    scale = ratio ** (-exp)
    if target_h_override and h_pt and target_h_override < h_pt * 0.45:
        scale = min(scale, 0.90)
    scale = min(1.0, max(scale, min_pt / o))
    out = max(min_pt, o * scale)
    fl = _ppt_readability_floor_pt(o, new_text, original_text, only_translation)
    return max(out, min(fl, o))


def _ppt_wrap_layout_params(
    new_text: str,
    w_pt: Optional[float],
    *,
    only_translation: bool,
) -> Tuple[float, float, float]:
    """与 _ppt_bilingual_fit_font_pt 一致的换行估算参数 (width_scale, height_pessimism, line_factor)。"""
    lr = _ppt_latin_letter_ratio(new_text or "")
    nt_for_len = (new_text or "").strip()
    if only_translation:
        width_scale = 0.87 if lr > 0.32 else (0.91 if lr > 0.16 else 0.94)
        height_pessimism = 1.11 if lr > 0.24 else 1.06
        line_factor = 1.30 if lr > 0.20 else 1.26
    else:
        width_scale = 0.88 if lr > 0.38 else (0.92 if lr > 0.22 else 0.96)
        height_pessimism = 1.12 if lr > 0.28 else 1.07
        line_factor = 1.32 if lr > 0.25 else 1.27
    if w_pt is not None and w_pt < 72 and len(nt_for_len) > 55:
        width_scale *= 0.94
    return width_scale, height_pessimism, line_factor


def _ppt_relax_unified_font_pt_for_shape(
    base_pt: float,
    pending: List[Dict[str, Any]],
    para_target_h: Dict[int, float],
    text_frame,
    shape,
    frame_orig_pt_ref: float,
    *,
    only_translation: bool = False,
    max_rel_ratio: float = 1.12,
    height_slack: float = 1.038,
) -> float:
    """
    多段已统一为 base_pt 后，在仍满足各段估算高度的情况下略增字号（利用框内留白）。
    """
    min_pt = 6.5
    try:
        b = float(base_pt)
        cap = float(frame_orig_pt_ref)
    except (TypeError, ValueError):
        return base_pt
    if b < min_pt or not pending:
        return base_pt
    w_pt, h_pt = _ppt_text_frame_inner_box_pt(text_frame, shape)
    if w_pt is None or h_pt is None or w_pt < 16 or h_pt < 14:
        return base_pt

    def _fits_at(cand: float) -> bool:
        for item in pending:
            th_raw = para_target_h.get(item["p_idx"])
            if th_raw is None or th_raw <= 0:
                continue
            target_h = max(min_pt * 2.0, float(th_raw))
            if only_translation:
                target_h *= 0.945
            allowed = target_h * height_slack
            ws, hp, lf = _ppt_wrap_layout_params(
                item["new_text"], w_pt, only_translation=only_translation
            )
            est = _ppt_wrapped_content_height_pt(
                item["new_text"],
                cand,
                w_pt,
                line_factor=lf,
                width_scale=ws,
                height_pessimism=hp,
            )
            if est > allowed:
                return False
        return True

    if not _fits_at(b):
        return base_pt

    hi = min(cap, b * max_rel_ratio)
    if hi <= b + 0.08:
        return b

    lo, best = b, b
    for _ in range(16):
        mid = (lo + hi) / 2.0
        if _fits_at(mid):
            best = mid
            lo = mid
        else:
            hi = mid
        if hi - lo < 0.055:
            break
    return max(b, best)


def copy_font_style(src_font, dest_font):
    """
    复制字体样式
    Args:
        src_font: 源字体对象 (Run.font)
        dest_font: 目标字体对象 (Run.font)
    """
    # 字体名称
    if src_font.name:
        dest_font.name = src_font.name
        # 针对中文字体的特殊处理 (Word/PPT通用)
        try:
            from docx.oxml.ns import qn
            if hasattr(dest_font, 'element'):
                dest_font.element.rPr.rFonts.set(qn('w:eastAsia'), src_font.name)
        except:
            pass
            
    # 字体大小
    if src_font.size:
        dest_font.size = src_font.size
    
    # 颜色
    try:
        if hasattr(src_font, 'color') and src_font.color and src_font.color.type:
            if src_font.color.type == 1: # RGB
                dest_font.color.rgb = src_font.color.rgb
            elif src_font.color.type == 2: # THEME
                dest_font.color.theme_color = src_font.color.theme_color
    except:
        pass

    # 粗体/斜体/下划线
    dest_font.bold = src_font.bold
    dest_font.italic = src_font.italic
    dest_font.underline = src_font.underline


def _ppt_font_effective_typeface(src_font) -> Optional[str]:
    """读取 PPT 运行字体：优先拉丁 typeface，否则 a:ea（中文常用）。"""
    if src_font is None:
        return None
    try:
        if src_font.name:
            return str(src_font.name)
        rPr = src_font._element
        from pptx.oxml.ns import qn

        ea = rPr.find(qn("a:ea"))
        if ea is not None:
            tf = ea.get("typeface")
            if tf:
                return str(tf)
    except Exception:
        pass
    return None


def copy_ppt_font_style(src_font, dest_font, fallback_name: str = "Microsoft YaHei") -> None:
    """复制 PPT 字符格式，并同步 a:latin 与 a:ea，避免中文显示成主题宋体与其它句不一致。"""
    if dest_font is None:
        return
    name = _ppt_font_effective_typeface(src_font) if src_font else None
    if not name:
        name = fallback_name
    try:
        from pptx.oxml import OxmlElement
        from pptx.oxml.ns import qn

        dest_font.name = name
        rPr = dest_font._element
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = OxmlElement("a:ea")
            ea.set("typeface", name)
            lat = rPr.find(qn("a:latin"))
            if lat is not None:
                lat.addnext(ea)
            else:
                rPr.append(ea)
        else:
            ea.set("typeface", name)

        if src_font:
            if src_font.size:
                dest_font.size = src_font.size
            if src_font.bold is not None:
                dest_font.bold = src_font.bold
            if src_font.italic is not None:
                dest_font.italic = src_font.italic
            if src_font.underline is not None:
                dest_font.underline = src_font.underline
            try:
                if (
                    hasattr(src_font, "color")
                    and src_font.color
                    and hasattr(src_font.color, "rgb")
                    and src_font.color.rgb
                ):
                    dest_font.color.rgb = src_font.color.rgb
            except Exception:
                pass
    except Exception:
        try:
            dest_font.name = fallback_name
        except Exception:
            pass


class FileGenerator:
    """文件生成器基类"""
    
    def __init__(
        self,
        file_path: str,
        translations: Dict[str, str],
        original_structure: Dict[str, Dict[str, Any]],
        source_lang: str = "en",
        target_lang: str = "zh-Hans",
        output_layout: str = "bilingual",
        glossary_file_path: str = "",
        industry_preset: str = "general",
    ):
        self.file_path = file_path
        self.translations = translations
        self.original_structure = original_structure
        self.source_lang = source_lang or "en"
        self.target_lang = target_lang or "zh-Hans"
        self.output_layout_for_name = output_layout or "bilingual"
        self.glossary_file_path = (glossary_file_path or "").strip()
        self.industry_preset_for_name = (industry_preset or "general").strip() or "general"
        self.output_path = self._generate_output_path()
    
    def _generate_output_path(self) -> str:
        """
        生成输出文件路径：
        原名_源语言_目标语言_翻译行业_版式_[词汇表主文件名]_XuYuan.ext
        翻译行业为界面选项对应的中文名（净化后）；版式为「双语对照」或「纯译文」；
        仅当词汇表路径有效且含至少一对词条时插入词汇表段。
        """
        return compute_translation_output_path(
            self.file_path,
            source_lang=self.source_lang,
            target_lang=self.target_lang,
            output_layout=self.output_layout_for_name,
            glossary_file_path=self.glossary_file_path,
            industry_preset=self.industry_preset_for_name,
        )
    
    def generate(self, format_type: str = "bilingual") -> str:
        """
        生成翻译后的文件
        Args:
            format_type: 输出版式："bilingual"（对照）或 "target_only"/"chinese_only"（仅译文）
        Returns:
            输出文件路径
        """
        raise NotImplementedError("子类必须实现generate方法")


class WordGenerator(FileGenerator):
    """Word文档生成器"""
    
    def generate(self, format_type: str = "bilingual") -> str:
        """生成翻译后的Word文档"""
        try:
            from docx import Document
            from docx.shared import Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml import OxmlElement
            
            doc = Document(self.file_path)
            
            # 创建编号到翻译的映射
            translation_map = {}
            for number, translation in self.translations.items():
                if number in self.original_structure:
                    key = self.original_structure[number]['key']
                    translation_map[key] = translation
            
            def _insert_paragraph_after_keep_order(paragraph):
                # 创建一个真正的段落对象，然后把它的 XML 节点移动到 paragraph 后面
                new_para = paragraph._parent.add_paragraph("")
                paragraph._p.addnext(new_para._p)
                return new_para

            def _copy_paragraph_format(src_para, dest_para):
                try:
                    dest_para.style = src_para.style
                except Exception:
                    pass
                try:
                    dest_para.alignment = src_para.alignment
                except Exception:
                    pass
                try:
                    spf = src_para.paragraph_format
                    dpf = dest_para.paragraph_format
                    dpf.left_indent = spf.left_indent
                    dpf.right_indent = spf.right_indent
                    dpf.first_line_indent = spf.first_line_indent
                    dpf.space_before = spf.space_before
                    dpf.space_after = spf.space_after
                    dpf.line_spacing = spf.line_spacing
                    dpf.keep_together = spf.keep_together
                    dpf.keep_with_next = spf.keep_with_next
                    dpf.page_break_before = spf.page_break_before
                    dpf.widow_control = spf.widow_control
                except Exception:
                    pass

            def _is_heading_paragraph(paragraph) -> bool:
                try:
                    style_name = (paragraph.style.name or "").strip().lower()
                except Exception:
                    style_name = ""
                # 常见：Heading 1/2/3... 或中文“标题 1/2/3...”
                if "heading" in style_name:
                    return True
                if "标题" in (paragraph.style.name or ""):
                    return True
                return False

            def apply_translation_to_paragraph(paragraph, key, in_table_cell: bool = False):
                if key in translation_map:
                    translation = translation_map[key]
                    original_text = paragraph.text
                    # 表格单元格提取的 cell.text 往往把多个段落用 '\n' 拼在一起，
                    # 直接写回到单一段落会导致 Word 生成“多行重复”的视觉效果。
                    # 这里仅在表格单元格场景下，把换行合并为空格，避免异常堆叠。
                    if in_table_cell and not is_target_only_layout(format_type):
                        translation = (
                            translation.replace("\r\n", " ")
                            .replace("\r", " ")
                            .replace("\n", " ")
                            .strip()
                        )
                    
                    # 1. 获取原始格式 (尝试从第一个非空Run获取)
                    src_run = None
                    if paragraph.runs:
                        src_run = paragraph.runs[0]
                    
                    if is_target_only_layout(format_type):
                        # 仅译文：直接替换当前段落文本。
                        # paragraph.clear() 只清除直接<w:r>元素，不触及深层文本框/图形。
                        paragraph.clear()
                        new_run = paragraph.add_run(translation)
                        if src_run:
                            copy_font_style(src_run.font, new_run.font)
                        else:
                            new_run.font.name = 'Microsoft YaHei'
                            new_run.font.size = Pt(10.5)
                    else:
                        # Heading/标题：不要插入新段落（否则可能继承标题样式，导致章节编号/目录层级异常）。
                        # 改为同段落追加译文，不换行。
                        if _is_heading_paragraph(paragraph):
                            # 保留原段落的编号/域/样式，只追加译文文本
                            sep = "  "
                            new_run = paragraph.add_run(f"{sep}{translation}")
                            if src_run:
                                copy_font_style(src_run.font, new_run.font)
                            else:
                                new_run.font.name = 'Microsoft YaHei'
                                new_run.font.size = Pt(10.5)
                            return

                        # 中英对照：保留原段落不动，译文作为“新段落”插在其后。
                        # 关键：避免在同一段落里加入手动换行(\n)，否则 Justify 会把空格拉伸得很夸张。
                        new_para = _insert_paragraph_after_keep_order(paragraph)
                        _copy_paragraph_format(paragraph, new_para)

                        # 如果原段落是两端对齐，译文段落改为左对齐，避免出现奇怪的字/词间距拉伸效果。
                        try:
                            if paragraph.alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                                new_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        except Exception:
                            pass

                        new_run = new_para.add_run(translation)
                        if src_run:
                            copy_font_style(src_run.font, new_run.font)
                        else:
                            new_run.font.name = 'Microsoft YaHei'
                            new_run.font.size = Pt(10.5)

            # 处理段落
            for para_idx, paragraph in enumerate(doc.paragraphs):
                key = f'paragraph_{para_idx}'
                apply_translation_to_paragraph(paragraph, key, in_table_cell=False)
            
            # ── 处理表格（含嵌套表格，按XML文档顺序） ──
            W_NS_FOR_TABLE = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            def _qnw(local_name):
                return f"{{{W_NS_FOR_TABLE}}}{local_name}"
            XML_NS_SPACE = 'http://www.w3.org/XML/1998/namespace'

            # 匹配解析器 _iter_tbl_row_elements / _iter_tr_cell_elements 中穿透 SDT 的逻辑
            def _xml_children_skip_sdt(parent):
                """遍历parent的直接子级，穿透<w:sdt>/<w:customXml>/<w:sdtContent>包装。"""
                for child in list(parent):
                    tag = child.tag
                    if tag in (_qnw('sdt'), _qnw('customXml'), _qnw('sdtContent')):
                        yield from _xml_children_skip_sdt(child)
                    else:
                        yield child

            def _apply_xml_para_translation(p_elem, translation):
                """在XML级别对<w:p>元素应用翻译。
                只移除直接子级的<w:r>run（它们是段落主体文字），
                保留深层的文本框/图形/超链接结构不受影响。
                """
                if is_target_only_layout(format_type):
                    for r in list(p_elem.findall(_qnw('r'))):
                        p_elem.remove(r)
                    new_r = OxmlElement('w:r')
                    new_t = OxmlElement('w:t')
                    new_t.text = translation
                    new_t.set(f'{{{XML_NS_SPACE}}}space', 'preserve')
                    new_r.append(new_t)
                    p_elem.append(new_r)
                else:
                    space_r = OxmlElement('w:r')
                    space_t = OxmlElement('w:t')
                    space_t.text = ' '
                    space_t.set(f'{{{XML_NS_SPACE}}}space', 'preserve')
                    space_r.append(space_t)
                    trans_r = OxmlElement('w:r')
                    trans_t = OxmlElement('w:t')
                    trans_t.text = translation
                    trans_t.set(f'{{{XML_NS_SPACE}}}space', 'preserve')
                    trans_r.append(trans_t)
                    p_elem.append(space_r)
                    p_elem.append(trans_r)

            def _process_cell_content(
                tc_elem, table_idx, row_idx, cell_idx, start_p_idx, para_xml_map
            ):
                """
                按XML文档顺序处理<w:tc>的直接子级（段落<w:p>和嵌套表格<w:tbl>），
                与解析器 _get_tc_all_texts 的遍历顺序完全一致。
                para_xml_map: {<w:p XML元素>: python-docx Paragraph对象} 映射表，嵌套单元格传None
                返回下一个可用p_idx。
                """
                current_p = start_p_idx
                for child in _xml_children_skip_sdt(tc_elem):
                    tag = child.tag
                    if tag == _qnw('p'):
                        # 解析器 _get_tc_all_texts 中跳过空段落（if text.strip()），生成器必须一致
                        has_extractable = False
                        for t_elem in child.iter(_qnw('t')):
                            if not t_elem.text or not t_elem.text.strip():
                                continue
                            # 跳过文本框内的文本
                            px = t_elem.getparent()
                            in_txbx = False
                            while px is not None and px is not child:
                                if px.tag == _qnw('txbxContent'):
                                    in_txbx = True
                                    break
                                px = px.getparent()
                            if not in_txbx:
                                has_extractable = True
                                break
                        if not has_extractable:
                            # 与解析器一致跳过此段落，不消耗p_idx
                            continue

                        key = f'table_{table_idx}_row_{row_idx}_cell_{cell_idx}_para_{current_p}'
                        # 用XML元素对比精确匹配python-docx Paragraph，而非依赖顺序
                        if para_xml_map is not None and child in para_xml_map:
                            p_obj = para_xml_map[child]
                            apply_translation_to_paragraph(p_obj, key, in_table_cell=True)
                        else:
                            # 嵌套单元格段落、或被SDT包裹但cell.paragraphs不包含的段落
                            if key in translation_map:
                                _apply_xml_para_translation(child, translation_map[key])
                        current_p += 1
                    elif tag == _qnw('tbl'):
                        # 嵌套表格：按行→单元格遍历（穿透SDT）
                        for tr in _xml_children_skip_sdt(child):
                            if tr.tag == _qnw('tr'):
                                for tc in _xml_children_skip_sdt(tr):
                                    if tc.tag == _qnw('tc'):
                                        # 嵌套单元格：所有段落走XML级别
                                        current_p = _process_cell_content(
                                            tc, table_idx, row_idx, cell_idx,
                                            current_p, None
                                        )
                return current_p

            for table_idx, table in enumerate(doc.tables):
                # 按XML<w:tc>枚举，与解析器 _iter_tr_cell_elements 一致（应对合并单元格）
                for row_idx, tr in enumerate(table.rows):
                    # 解析器迭代XML tr的子级，我们需要同样的tc_list
                    row_xml = tr._tr
                    tc_list = []
                    for xml_child in _xml_children_skip_sdt(row_xml):
                        if xml_child.tag == _qnw('tc'):
                            tc_list.append(xml_child)
                    for cell_xml_idx, tc_elem in enumerate(tc_list):
                        # 找python-docx中第一个指向此_tc的Cell（用于获取paragraphs）
                        cell_for_paras = None
                        for docx_cell in tr.cells:
                            if docx_cell._tc is tc_elem:
                                cell_for_paras = docx_cell
                                break
                        para_xml_map = {}
                        if cell_for_paras is not None:
                            para_xml_map = {p._element: p for p in cell_for_paras.paragraphs}
                        _process_cell_content(
                            tc_elem, table_idx, row_idx, cell_xml_idx,
                            0, para_xml_map
                        )
            
            doc.save(self.output_path)
            logger.info(f"Word文档已保存到: {self.output_path}")
            return self.output_path
            
        except ImportError:
            raise FileParseError("请安装python-docx库")
        except Exception as e:
            raise FileParseError(f"生成Word文档失败: {str(e)}")


class PowerPointGenerator(FileGenerator):
    """PowerPoint文档生成器 (优化版：修复软回车和组合图形漏翻问题)"""
    
    def generate(self, format_type: str = "bilingual") -> str:
        """生成翻译后的PowerPoint文档"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.util import Pt

            prs = Presentation(self.file_path)
            
            # 1. 构建形状层映射表（排除 OOXML 条目，避免 slide_index=None 与 (None,text) 冲突）
            text_content_map = {}
            for number, data in self.original_structure.items():
                if data.get("type") == "ppt_ooxml_text":
                    continue
                if data.get("slide_index") is None:
                    continue
                if "text" in data and number in self.translations:
                    clean_key_text = data["text"].strip()
                    k = (data["slide_index"], clean_key_text)
                    text_content_map[k] = self.translations[number]

            total_shapes = 0
            translated_shapes = 0
            
            # 递归遍历函数
            def _iter_shapes(shapes):
                for shape in shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from _iter_shapes(shape.shapes)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                yield cell
                    else:
                        yield shape

            # 处理幻灯片
            for slide_idx, slide in enumerate(prs.slides):
                for shape in _iter_shapes(slide.shapes):
                    
                    text_frame = None
                    try:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            text_frame = shape.text_frame
                        elif hasattr(shape, "text_frame"):
                            text_frame = shape.text_frame
                    except Exception:
                        continue 
                        
                    if not text_frame:
                        continue
                        
                    if not text_frame.text.strip():
                        continue
                        
                    total_shapes += 1
                    shape_modified = False

                    wide_body = _ppt_shape_is_wide_body_text_block(text_frame, shape)
                    # 宽幅正文：多段均分高度 + 统一字号；窄框/图示内文字：沿用加权高度 + 逐段字号，避免叠字
                    if wide_body:
                        para_target_h = _ppt_equal_paragraph_target_heights_pt(
                            text_frame, shape
                        )
                        if not para_target_h:
                            para_target_h = _ppt_paragraph_target_heights_pt(
                                text_frame, shape
                            )
                    else:
                        para_target_h = _ppt_paragraph_target_heights_pt(
                            text_frame, shape
                        )
                    frame_orig_pt_ref = _ppt_frame_max_run_font_pt(text_frame)

                    # 勿用 TEXT_TO_FIT_SHAPE：与手动计算字号冲突，易出现仍超范围、叠图
                    if text_frame.auto_size is None:
                        try:
                            text_frame.auto_size = MSO_AUTO_SIZE.NONE
                            text_frame.word_wrap = True
                        except Exception:
                            pass

                    only_tgt = is_target_only_layout(format_type)
                    pending: List[Dict[str, Any]] = []

                    # 先收集本形状内需写入的段落，再对多段统一字号（取各段适配值的最小值）
                    for p_idx, paragraph in enumerate(list(text_frame.paragraphs)):
                        original_text = paragraph.text
                        stripped_text = original_text.strip()
                        if not stripped_text:
                            continue

                        found_translation = None
                        used_full_match = False

                        key_full = (slide_idx, stripped_text)
                        if key_full in text_content_map:
                            found_translation = text_content_map[key_full]
                            used_full_match = True

                        separator = "\x0b" if "\x0b" in original_text else None
                        segments = []
                        translated_segments = []

                        if used_full_match:
                            segments = [original_text]
                            if only_tgt:
                                translated_segments = [found_translation]
                            else:
                                translated_segments = [
                                    f"{original_text} {found_translation}"
                                ]
                        else:
                            if separator:
                                segments = original_text.split(separator)
                            else:
                                segments = [original_text]

                            for seg in segments:
                                seg_clean = seg.strip()
                                if not seg_clean:
                                    translated_segments.append(seg)
                                    continue

                                key = (slide_idx, seg_clean)
                                trans = text_content_map.get(key)
                                if not trans:
                                    key_fallback = (
                                        slide_idx,
                                        seg_clean.replace("\xa0", " "),
                                    )
                                    trans = text_content_map.get(key_fallback)

                                if trans:
                                    if only_tgt:
                                        translated_segments.append(trans)
                                    else:
                                        translated_segments.append(f"{seg} {trans}")
                                else:
                                    translated_segments.append(seg)

                        is_modified = used_full_match or any(
                            s != o for s, o in zip(translated_segments, segments)
                        )

                        if not is_modified:
                            continue

                        join_char = separator if separator and not used_full_match else ""
                        if used_full_match:
                            new_text = translated_segments[0]
                        else:
                            new_text = join_char.join(translated_segments)

                        src_font = paragraph.runs[0].font if paragraph.runs else None
                        try:
                            orig_pt = (
                                float(src_font.size.pt)
                                if src_font and src_font.size
                                else 12.0
                            )
                        except Exception:
                            orig_pt = 12.0

                        pending.append(
                            {
                                "paragraph": paragraph,
                                "original_text": original_text,
                                "new_text": new_text,
                                "src_font": src_font,
                                "orig_pt": orig_pt,
                                "p_idx": p_idx,
                            }
                        )

                    apply_pt: Optional[float] = None
                    if len(pending) > 1 and wide_body:
                        fit_pts: List[float] = []
                        for item in pending:
                            fit_pts.append(
                                _ppt_bilingual_fit_font_pt(
                                    frame_orig_pt_ref,
                                    item["original_text"],
                                    item["new_text"],
                                    text_frame,
                                    shape,
                                    target_h_override=para_target_h.get(
                                        item["p_idx"]
                                    ),
                                    only_translation=only_tgt,
                                )
                            )
                        apply_pt = min(fit_pts) if fit_pts else None
                        if apply_pt is not None and pending:
                            apply_pt = _ppt_relax_unified_font_pt_for_shape(
                                apply_pt,
                                pending,
                                para_target_h,
                                text_frame,
                                shape,
                                frame_orig_pt_ref,
                                only_translation=only_tgt,
                            )
                    elif len(pending) > 1 and not wide_body:
                        for item in pending:
                            item["fit_pt"] = _ppt_bilingual_fit_font_pt(
                                item["orig_pt"],
                                item["original_text"],
                                item["new_text"],
                                text_frame,
                                shape,
                                target_h_override=para_target_h.get(
                                    item["p_idx"]
                                ),
                                only_translation=only_tgt,
                            )
                    elif len(pending) == 1:
                        one = pending[0]
                        apply_pt = _ppt_bilingual_fit_font_pt(
                            one["orig_pt"],
                            one["original_text"],
                            one["new_text"],
                            text_frame,
                            shape,
                            target_h_override=para_target_h.get(one["p_idx"]),
                            only_translation=only_tgt,
                        )

                    for item in pending:
                        paragraph = item["paragraph"]
                        new_text = item["new_text"]
                        src_font = item["src_font"]

                        paragraph.clear()
                        new_run = paragraph.add_run()
                        new_run.text = new_text

                        if src_font:
                            copy_ppt_font_style(src_font, new_run.font)
                        else:
                            copy_ppt_font_style(None, new_run.font)
                            try:
                                new_run.font.size = Pt(12)
                            except Exception:
                                pass

                        pt_use = apply_pt
                        if pt_use is None and item.get("fit_pt") is not None:
                            pt_use = item["fit_pt"]
                        if pt_use is not None:
                            try:
                                new_run.font.size = Pt(pt_use)
                            except Exception:
                                pass

                        shape_modified = True

                    if shape_modified:
                        try:
                            text_frame.auto_size = MSO_AUTO_SIZE.NONE
                            text_frame.word_wrap = True
                        except Exception:
                            pass
                        
                    if shape_modified:
                        translated_shapes += 1

            logger.info(
                f"PPT生成统计: 扫描到含文本形状 {total_shapes}, 成功翻译写入形状层 {translated_shapes}"
            )

            prs.save(self.output_path)

            # SmartArt / 图表 / 备注：在保存后按 Zip 内部件路径与节点序号写回
            from core.ppt_ooxml import apply_ppt_ooxml_translations

            only_ooxml = is_target_only_layout(format_type)
            ooxml_updates: List[Tuple[str, int, str]] = []
            for number, data in self.original_structure.items():
                if data.get("type") != "ppt_ooxml_text":
                    continue
                if number not in self.translations:
                    continue
                part_path = data["part_path"]
                t_idx = int(data["t_index"])
                orig = (data.get("text") or "").strip()
                trans = (self.translations.get(number) or "").strip()
                if only_ooxml:
                    new_text = trans if trans else orig
                else:
                    new_text = f"{orig} {trans}".strip() if trans else orig
                ooxml_updates.append((part_path, t_idx, new_text))

            ooxml_written = 0
            if ooxml_updates:
                try:
                    ooxml_written = apply_ppt_ooxml_translations(
                        self.output_path, ooxml_updates
                    )
                except Exception as ex:
                    logger.warning("PPT OOXML 补丁写入失败: %s", ex)

            logger.info(
                "PPT OOXML 补丁: 计划 %s 处, 实际写入文本节点 %s 个",
                len(ooxml_updates),
                ooxml_written,
            )

            ooxml_trans_count = sum(
                1
                for n, d in self.original_structure.items()
                if d.get("type") == "ppt_ooxml_text" and n in self.translations
            )
            slack = len(self.translations) - translated_shapes - ooxml_trans_count
            if slack > 50:
                logger.warning(
                    "注意: 约有 %s 条译文字段未在形状层或 OOXML 补丁中消耗，可能含母版或其它结构",
                    slack,
                )

            return self.output_path
            
        except ImportError:
            raise FileParseError("请安装python-pptx")
        except Exception as e:
            import traceback
            logger.error(traceback.format_exc())
            raise FileParseError(f"PPT生成失败: {str(e)}")


class ExcelGenerator(FileGenerator):
    """Excel文档生成器"""
    
    def generate(self, format_type: str = "bilingual") -> str:
        """生成翻译后的Excel文档"""
        try:
            from openpyxl import load_workbook
            from openpyxl.comments import Comment
            
            wb = load_workbook(self.file_path)
            
            # 创建编号到翻译的映射
            translation_map = {}
            for number, translation in self.translations.items():
                if number in self.original_structure:
                    key = self.original_structure[number]['key']
                    translation_map[key] = translation
            
            # 处理单元格
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue

                        key = f'sheet_{sheet_name}_cell_{cell.coordinate}'
                        if key not in translation_map:
                            continue

                        translation = translation_map[key]
                        original_text = str(cell.value)

                        # 公式单元格：不改写 value，避免破坏公式；把译文写入批注。
                        if cell.data_type == 'f':
                            # 保留已有批注（若存在）并追加“译文”段。
                            existing = ""
                            if cell.comment and cell.comment.text:
                                existing = str(cell.comment.text).strip()

                            trans_block = f"译文：{translation}" if translation else "译文："
                            if existing:
                                if trans_block in existing:
                                    comment_text = existing
                                else:
                                    comment_text = f"{existing}\n\n{trans_block}"
                            else:
                                comment_text = trans_block

                            # openpyxl 的 Comment 需要 author；给定固定作者避免为空。
                            cell.comment = Comment(comment_text, "XuYuan")
                            continue

                        if is_target_only_layout(format_type):
                            cell.value = translation
                        else:  # bilingual
                            # 文本以 '=' 开头时，避免被 Excel 误判为公式
                            bilingual_text = f"{original_text}\n{translation}"
                            if bilingual_text.startswith('='):
                                bilingual_text = "'" + bilingual_text
                            cell.value = bilingual_text
            
            # 保存文件
            wb.save(self.output_path)
            logger.info(f"Excel文档已保存到: {self.output_path}")
            return self.output_path
            
        except ImportError:
            raise FileParseError("请安装openpyxl库: pip install openpyxl")
        except Exception as e:
            raise FileParseError(f"生成Excel文档失败: {str(e)}")


class PDFGenerator(FileGenerator):
    """原地擦除+写回 PDF：纯中文保持原字号；中英对照将字号减小并上下分区。"""

    # 使用内置 CJK 字体，避免外部 fontfile 在部分系统/中文抽取中不稳定的问题。
    # 注意：这不保证“复制出来的中文完全正确”，但至少能稳定渲染并可见覆盖原文。
    _FITZ_CN_FONTNAME = "china-ss"
    _FITZ_EN_FONTNAME = "helv"

    @staticmethod
    def _wrap_text_greedy(
        text: str,
        fontname: str,
        fontsize: float,
        max_width: float,
    ) -> List[str]:
        """按宽度贪心换行。支持 '\n' 强制换行。"""
        if not text:
            return []

        text = str(text).replace("\r\n", "\n").replace("\r", "\n")
        lines: List[str] = []

        def measure(s: str) -> float:
            import fitz

            return float(fitz.get_text_length(s, fontname=fontname, fontsize=fontsize))

        for para in text.split("\n"):
            para = para.strip()
            if not para:
                lines.append("")
                continue

            # 如果有空格，按“词”来换行；否则逐字符换行（中文通常没有空格）
            if " " in para:
                tokens = [t for t in para.split(" ") if t]
                cur = ""
                for tok in tokens:
                    cand = tok if not cur else cur + " " + tok
                    if measure(cand) <= max_width:
                        cur = cand
                    else:
                        if cur:
                            lines.append(cur)
                        cur = tok
                if cur:
                    lines.append(cur)
            else:
                cur = ""
                for ch in para:
                    cand = ch if not cur else cur + ch
                    if measure(cand) <= max_width:
                        cur = cand
                    else:
                        if cur:
                            lines.append(cur)
                        cur = ch
                if cur:
                    lines.append(cur)

        return lines

    @staticmethod
    def _insert_wrapped_text(
        page,
        rect,
        text: str,
        font_size: float,
        fontname: str,
        color=None,
        min_font_size: float = 2.6,
        dry_run: bool = False,
    ) -> float:
        """把 text 写入 rect，必要时自动缩小字号。返回最终字号；返回 <min 则表示未写入。"""
        import fitz

        if not text:
            return min_font_size - 1

        rr = fitz.Rect(rect)
        x0, y0, x1, y1 = rr.x0, rr.y0, rr.x1, rr.y1
        width = max(0.0, x1 - x0)
        height = max(0.0, y1 - y0)
        if width <= 1 or height <= 1:
            return min_font_size - 1

        pad_x = min(2.0, max(0.5, font_size * 0.05))
        usable_w = max(1.0, width - 2 * pad_x)

        fs = float(font_size)
        while fs >= float(min_font_size):
            line_height = fs * 1.25
            if line_height <= 0:
                fs *= 0.9
                continue
            max_lines = int(height // line_height)
            if max_lines <= 0:
                fs *= 0.9
                continue

            lines = PDFGenerator._wrap_text_greedy(
                text=text,
                fontname=fontname,
                fontsize=fs,
                max_width=usable_w,
            )

            # strip 行尾空格，但保留空行（由 '\n\n' 之类产生）
            if len(lines) <= max_lines and all(
                float(fitz.get_text_length(line, fontname=fontname, fontsize=fs)) <= usable_w + 0.01
                for line in lines
            ):
                # fitz 的 insert_text y 是 baseline（以页面左上为原点、向下为正方向）
                cur_baseline = y0 + fs
                for i, line in enumerate(lines):
                    if i >= max_lines:
                        break
                    if not dry_run:
                        page.insert_text(
                            (x0 + pad_x, cur_baseline + i * line_height),
                            line,
                            fontsize=fs,
                            fontname=fontname,
                            color=color,
                            render_mode=0,
                        )
                return fs

            fs *= 0.9

        return min_font_size - 1

    @staticmethod
    def _normalize_pdf_color(raw_color):
        """
        将 pdfplumber 的颜色值归一化为 PyMuPDF 颜色：
        - 灰度 float/int -> (g,g,g)
        - RGB tuple/list -> (r,g,b)
        分量统一到 0..1；无法识别则返回 None（走默认黑色）。
        """
        if raw_color is None:
            return None

        def norm(v):
            try:
                fv = float(v)
            except Exception:
                return None
            if fv < 0:
                fv = 0.0
            if fv > 1:
                # 兼容 0..255
                fv = fv / 255.0
            return max(0.0, min(1.0, fv))

        if isinstance(raw_color, (int, float)):
            g = norm(raw_color)
            return (g, g, g) if g is not None else None

        if isinstance(raw_color, (list, tuple)):
            vals = [norm(v) for v in raw_color]
            vals = [v for v in vals if v is not None]
            if not vals:
                return None
            if len(vals) == 1:
                return (vals[0], vals[0], vals[0])
            # RGB 或 RGBA（取前三个）
            return tuple(vals[:3]) if len(vals) >= 3 else (vals[0], vals[0], vals[0])

        return None

    def generate(self, format_type: str = "bilingual") -> str:
        """原地擦除原文并写回译文到 PDF 指定区域。"""
        try:
            import fitz

            from collections import defaultdict

            segs_by_page: Dict[int, List[Tuple[str, Dict[str, Any]]]] = defaultdict(list)
            for number, data in self.original_structure.items():
                page_idx = int(data.get("page_index", 0) or 0)
                rect = data.get("rect")
                if not rect or len(rect) != 4:
                    continue
                segs_by_page[page_idx].append((number, data))

            doc = fitz.open(self.file_path)

            for page_idx, segs in sorted(segs_by_page.items()):
                page = doc[page_idx]
                segs.sort(key=lambda x: int(x[1].get("paragraph_index", 0) or 0))

                # 先规划“可写回动作”：只有能写回的块才允许被擦除
                min_fit_size = 2.6
                actions = []
                for number, data in segs:
                    rect = data["rect"]
                    x0, y0, x1, y1 = map(float, rect)
                    base_size = float(data.get("font_size", 10.0) or 10.0)

                    orig = data.get("text", "")
                    # 如果该编号未返回翻译，回退到原文，避免页面出现占位符
                    trans = self.translations.get(number) or orig
                    text_color = self._normalize_pdf_color(data.get("text_color"))

                    rr = fitz.Rect(x0, y0, x1, y1)

                    if is_target_only_layout(format_type):
                        # 纯中文优先保持原字号，仅在放不下时自动缩小
                        font_size = max(6.0, base_size)
                        fit_fs = self._insert_wrapped_text(
                            page,
                            rr,
                            trans,
                            font_size=font_size,
                            fontname=self._FITZ_CN_FONTNAME,
                            color=text_color,
                            min_font_size=min_fit_size,
                            dry_run=True,
                        )
                        if fit_fs >= min_fit_size:
                            actions.append(
                                {
                                    "rect": rr,
                                    "writes": [
                                        {
                                            "rect": rr,
                                            "text": trans,
                                            "font_size": font_size,
                                            "fontname": self._FITZ_CN_FONTNAME,
                                            "color": text_color,
                                            "min_font_size": min_fit_size,
                                        }
                                    ],
                                }
                            )
                    else:
                        # 上英文 + 下中文，先用更自然的初始字号，再按区域自动缩放
                        # 双语初始字号略放大，再由自动缩放兜底
                        font_size = max(6.0, base_size * 0.9)
                        mid_y = rr.y0 + rr.height / 2.0
                        top_rr = fitz.Rect(rr.x0, rr.y0, rr.x1, mid_y)
                        bottom_rr = fitz.Rect(rr.x0, mid_y, rr.x1, rr.y1)

                        fit_top = self._insert_wrapped_text(
                            page,
                            top_rr,
                            orig,
                            font_size=font_size,
                            fontname=self._FITZ_EN_FONTNAME,
                            color=text_color,
                            min_font_size=min_fit_size,
                            dry_run=True,
                        )
                        fit_bottom = self._insert_wrapped_text(
                            page,
                            bottom_rr,
                            trans,
                            font_size=font_size,
                            fontname=self._FITZ_CN_FONTNAME,
                            color=text_color,
                            min_font_size=min_fit_size,
                            dry_run=True,
                        )

                        if fit_top >= min_fit_size and fit_bottom >= min_fit_size:
                            actions.append(
                                {
                                    "rect": rr,
                                    "writes": [
                                        {
                                            "rect": top_rr,
                                            "text": orig,
                                            "font_size": font_size,
                                            "fontname": self._FITZ_EN_FONTNAME,
                                            "color": text_color,
                                            "min_font_size": min_fit_size,
                                        },
                                        {
                                            "rect": bottom_rr,
                                            "text": trans,
                                            "font_size": font_size,
                                            "fontname": self._FITZ_CN_FONTNAME,
                                            "color": text_color,
                                            "min_font_size": min_fit_size,
                                        },
                                    ],
                                }
                            )
                        else:
                            # 双语放不下时继续尝试：整块双语（原文+译文）进一步缩字号
                            bi_full_text = f"{orig}\n{trans}"
                            fit_bi_full = self._insert_wrapped_text(
                                page,
                                rr,
                                bi_full_text,
                                font_size=max(6.0, base_size * 0.9),
                                fontname=self._FITZ_CN_FONTNAME,
                                color=text_color,
                                min_font_size=min_fit_size,
                                dry_run=True,
                            )
                            if fit_bi_full >= min_fit_size:
                                actions.append(
                                    {
                                        "rect": rr,
                                        "writes": [
                                            {
                                                "rect": rr,
                                                "text": bi_full_text,
                                                "font_size": max(6.0, base_size * 0.9),
                                                "fontname": self._FITZ_CN_FONTNAME,
                                                "color": text_color,
                                                "min_font_size": min_fit_size,
                                            }
                                        ],
                                    }
                                )
                                continue

                            # 仍放不下时再降级：整块仅写中文；再不行才保留原文不擦除
                            fit_cn_full = self._insert_wrapped_text(
                                page,
                                rr,
                                trans,
                                font_size=max(6.0, base_size),
                                fontname=self._FITZ_CN_FONTNAME,
                                color=text_color,
                                min_font_size=min_fit_size,
                                dry_run=True,
                            )
                            if fit_cn_full >= min_fit_size:
                                actions.append(
                                    {
                                        "rect": rr,
                                        "writes": [
                                            {
                                                "rect": rr,
                                                "text": trans,
                                                "font_size": max(6.0, base_size),
                                                "fontname": self._FITZ_CN_FONTNAME,
                                                "color": text_color,
                                                "min_font_size": min_fit_size,
                                            }
                                        ],
                                    }
                                )

                # 只对“可写回动作”做擦除，避免出现内容被擦掉但未写回
                for act in actions:
                    page.add_redact_annot(act["rect"], fill=None, cross_out=False)

                # 只移除文本，避免图形和图片被“抹白”
                if actions:
                    page.apply_redactions(images=0, graphics=0, text=0)

                # 按计划执行写回
                for act in actions:
                    for w in act["writes"]:
                        self._insert_wrapped_text(
                            page,
                            w["rect"],
                            w["text"],
                            font_size=w["font_size"],
                            fontname=w["fontname"],
                            color=w["color"],
                            min_font_size=w.get("min_font_size", min_fit_size),
                            dry_run=False,
                        )

            doc.save(self.output_path)
            doc.close()
            logger.info(f"PDF 翻译结果已保存: {self.output_path}")
            return self.output_path

        except FileParseError:
            raise
        except ImportError as e:
            raise FileParseError("请安装 PyMuPDF: pip install PyMuPDF") from e
        except Exception as e:
            raise FileParseError(f"生成PDF翻译结果失败: {str(e)}")


def create_generator(
    file_path: str,
    translations: Dict[str, str],
    original_structure: Dict[str, Dict[str, Any]],
    source_lang: str = "en",
    target_lang: str = "zh-Hans",
    output_layout: str = "bilingual",
    glossary_file_path: str = "",
    industry_preset: str = "general",
) -> FileGenerator:
    """根据文件类型创建相应的生成器"""
    file_ext = Path(file_path).suffix.lower()
    
    generator_map = {
        # Word
        '.docx': WordGenerator,
        '.docm': WordGenerator, # 新增
        
        # PPT
        '.pptx': PowerPointGenerator,
        '.pptm': PowerPointGenerator, # 新增
        
        # Excel
        '.xlsx': ExcelGenerator,
        '.xlsm': ExcelGenerator, # 新增
        
        # PDF
        '.pdf': PDFGenerator,
    }
    
    generator_class = generator_map.get(file_ext)
    if generator_class is None:
        raise FileParseError(f"不支持生成的文件类型: {file_ext}")
    
    return generator_class(
        file_path,
        translations,
        original_structure,
        source_lang=source_lang,
        target_lang=target_lang,
        output_layout=output_layout,
        glossary_file_path=glossary_file_path,
        industry_preset=industry_preset,
    )