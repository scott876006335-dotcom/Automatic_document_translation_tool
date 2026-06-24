"""文件解析模块"""
import os
from typing import Dict, List, Tuple, Any
from pathlib import Path

from utils.exceptions import FileParseError
from utils.logger import setup_logger

logger = setup_logger()


class FileParser:
    """文件解析器基类"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        if not os.path.exists(file_path):
            raise FileParseError(f"文件不存在: {file_path}")
    
    def parse(self) -> Dict[str, Dict[str, Any]]:
        """解析文件，返回编号文本字典"""
        raise NotImplementedError("子类必须实现parse方法")
    
    def get_file_type(self) -> str:
        """获取文件类型"""
        return Path(self.file_path).suffix.lower()


class WordParser(FileParser):
    """Word文档解析器"""

    # Word XML 命名空间
    _W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    @classmethod
    def _qn(cls, local_name: str) -> str:
        return f"{{{cls._W_NS}}}{local_name}"

    # ----------------------------------------------------------------
    #  body 级别：穿透 SDT / customXml，找到所有顶层 <w:p> 和 <w:tbl>
    # ----------------------------------------------------------------
    @classmethod
    def _iter_body_block_elements(cls, body):
        """
        从 <w:body> 中按文档顺序产出所有块级元素，
        穿透 <w:sdt> / <w:customXml> / <w:sdtContent>。
        Yields: ('p', element)  或  ('tbl', element)
        """
        tag_p   = cls._qn("p")
        tag_tbl = cls._qn("tbl")
        tag_sdt = cls._qn("sdt")
        tag_customXml  = cls._qn("customXml")
        tag_sdtContent = cls._qn("sdtContent")

        def _collect(parent):
            for child in parent:
                tag = child.tag
                if tag == tag_p:
                    yield ("p", child)
                elif tag == tag_tbl:
                    yield ("tbl", child)
                elif tag in (tag_sdt, tag_customXml, tag_sdtContent):
                    yield from _collect(child)
                # 其他标签（如 w:sectPr）忽略

        yield from _collect(body)

    # ----------------------------------------------------------------
    #  table 内部：穿透 SDT 找 <w:tr>
    # ----------------------------------------------------------------
    @classmethod
    def _iter_tbl_row_elements(cls, tbl_element):
        tag_tr  = cls._qn("tr")
        tag_tbl = cls._qn("tbl")
        tag_sdt = cls._qn("sdt")
        tag_customXml  = cls._qn("customXml")
        tag_sdtContent = cls._qn("sdtContent")

        def _collect(parent):
            for child in parent:
                tag = child.tag
                if tag == tag_tr:
                    yield child
                elif tag in (tag_sdt, tag_customXml, tag_sdtContent):
                    yield from _collect(child)
                # 遇到嵌套 <w:tbl> 则跳过

        yield from _collect(tbl_element)

    # ----------------------------------------------------------------
    #  row 内部：穿透 SDT 找 <w:tc>
    # ----------------------------------------------------------------
    @classmethod
    def _iter_tr_cell_elements(cls, tr_element):
        tag_tc  = cls._qn("tc")
        tag_sdt = cls._qn("sdt")
        tag_customXml  = cls._qn("customXml")
        tag_sdtContent = cls._qn("sdtContent")

        def _collect(parent):
            for child in parent:
                tag = child.tag
                if tag == tag_tc:
                    yield child
                elif tag in (tag_sdt, tag_customXml, tag_sdtContent):
                    yield from _collect(child)

        yield from _collect(tr_element)

    # ----------------------------------------------------------------
    #  cell 内部：提取段落文本，同时递归处理嵌套表格
    # ----------------------------------------------------------------
    @classmethod
    def _get_tc_all_texts(cls, tc_element) -> "List[str]":
        """
        从 <w:tc> 中提取所有纯文本段落，包括：
        - 直属 <w:p>（穿透 SDT 包裹）
        - 嵌套 <w:tbl> 中的所有单元格文本（递归）
        每个非空段落作为列表中一个元素返回。
        注意：排除文本框 <w:txbxContent> 内部的文本，避免混入单元格文本序列。
        """
        tag_p   = cls._qn("p")
        tag_t   = cls._qn("t")
        tag_tbl = cls._qn("tbl")
        tag_sdt = cls._qn("sdt")
        tag_customXml  = cls._qn("customXml")
        tag_sdtContent = cls._qn("sdtContent")
        tag_txbxContent = cls._qn("txbxContent")

        def _extract_p_text(p_element) -> str:
            """从 <w:p> 提取纯文本，跳过文本框内的 <w:t>。"""
            parts = []
            for t_elem in p_element.iter(tag_t):
                if not t_elem.text:
                    continue
                # 跳过文本框 <w:txbxContent> 内部的文本
                parent_iter = t_elem.getparent()
                in_txbx = False
                while parent_iter is not None and parent_iter is not p_element:
                    if parent_iter.tag == tag_txbxContent:
                        in_txbx = True
                        break
                    parent_iter = parent_iter.getparent()
                if not in_txbx:
                    parts.append(t_elem.text)
            return "".join(parts)

        results: "List[str]" = []

        def _collect(parent):
            for child in parent:
                tag = child.tag
                if tag == tag_p:
                    text = _extract_p_text(child)
                    if text.strip():
                        results.append(text)
                elif tag == tag_tbl:
                    # ★ 递归处理嵌套表格
                    for tr in cls._iter_tbl_row_elements(child):
                        for tc in cls._iter_tr_cell_elements(tr):
                            nested = cls._get_tc_all_texts(tc)
                            results.extend(nested)
                elif tag in (tag_sdt, tag_customXml, tag_sdtContent):
                    _collect(child)
                # 其他标签（w:tcPr 等）忽略

        _collect(tc_element)
        return results

    # ----------------------------------------------------------------
    #  从 <w:p> XML 元素提取纯文本（用于 body 级别段落）
    # ----------------------------------------------------------------
    @classmethod
    def _get_paragraph_text(cls, p_element) -> str:
        tag_t   = cls._qn("t")
        tag_tab = cls._qn("tab")
        tag_br  = cls._qn("br")
        tag_cr  = cls._qn("cr")
        tag_txbxContent = cls._qn("txbxContent")

        parts = []
        for n in p_element.iter():
            if n.tag == tag_t and n.text:
                # 跳过文本框(<w:txbxContent>)内部的文本
                parent_iter = n.getparent()
                in_txbx = False
                while parent_iter is not None and parent_iter is not p_element:
                    if parent_iter.tag == tag_txbxContent:
                        in_txbx = True
                        break
                    parent_iter = parent_iter.getparent()
                if not in_txbx:
                    parts.append(n.text)
            elif n.tag == tag_tab:
                parts.append("\t")
            elif n.tag in (tag_br, tag_cr):
                parts.append("\n")
        return "".join(parts)

    # ================================================================
    #  主解析方法
    # ================================================================
    def parse(self) -> "Dict[str, Dict[str, Any]]":
        """解析Word文档（最小可替换版：优先保证不漏文本）"""
        try:
            from docx import Document

            doc = Document(self.file_path)
            numbered_texts: Dict[str, Dict[str, Any]] = {}
            index = 1

            body = doc.element.body
            #seen_tc_set: set = set()

            para_counter = 0
            table_counter = 0

            # 调试计数
            body_p_count = 0
            table_cell_para_count = 0

            # 本地兜底：防止类里还没定义 _iter_body_block_elements
            def _iter_body_block_elements_fallback(body_element):
                tag_p = self._qn("p")
                tag_tbl = self._qn("tbl")
                tag_sdt = self._qn("sdt")
                tag_customXml = self._qn("customXml")
                tag_sdtContent = self._qn("sdtContent")

                def _collect(parent):
                    for child in parent:
                        tag = child.tag
                        if tag == tag_p:
                            yield ("p", child)
                        elif tag == tag_tbl:
                            yield ("tbl", child)
                        elif tag in (tag_sdt, tag_customXml, tag_sdtContent):
                            yield from _collect(child)

                yield from _collect(body_element)

            iter_body_blocks = (
                self._iter_body_block_elements(body)
                if hasattr(self, "_iter_body_block_elements")
                else _iter_body_block_elements_fallback(body)
            )

            for elem_type, elem in iter_body_blocks:
                # ========== 正文段落 ==========
                if elem_type == "p":
                    text = self._get_paragraph_text(elem).strip()
                    if text:
                        number = f"P{index:04d}"
                        numbered_texts[number] = {
                            "key": f"paragraph_{para_counter}",
                            "text": text,
                            "type": "paragraph",
                            "index": para_counter,
                            "style": {},  # 最小版先不取样式，确保不漏文本
                        }
                        index += 1
                        body_p_count += 1
                    para_counter += 1

                # ========== 表格 ==========
                elif elem_type == "tbl":
                    tbl_elem = elem
                    for row_idx, tr in enumerate(self._iter_tbl_row_elements(tbl_elem)):
                        for cell_idx, tc in enumerate(self._iter_tr_cell_elements(tr)):
                            # 去重：同一 tc XML 节点只处理一次
                            #tc_id = id(tc)
                            #if tc_id in seen_tc_set:
                            #    continue
                            #seen_tc_set.add(tc_id)

                            # 优先用你类里的 _get_tc_all_texts（支持嵌套表格）
                            if hasattr(self, "_get_tc_all_texts"):
                                para_texts = self._get_tc_all_texts(tc)
                            else:
                                # 兜底：旧函数名
                                para_texts = self._get_tc_paragraph_texts(tc)

                            for p_idx, p_text in enumerate(para_texts):
                                text = (p_text or "").strip()
                                if not text:
                                    continue

                                number = f"P{index:04d}"
                                numbered_texts[number] = {
                                    "key": f"table_{table_counter}_row_{row_idx}_cell_{cell_idx}_para_{p_idx}",
                                    "text": text,
                                    "type": "table_cell_paragraph",
                                    "table_index": table_counter,
                                    "row_index": row_idx,
                                    "cell_index": cell_idx,
                                    "paragraph_index": p_idx,
                                }
                                index += 1
                                table_cell_para_count += 1

                    table_counter += 1

            logger.info(
                f"Word文档解析完成，共提取 {len(numbered_texts)} 个文本段落 "
                f"(body={body_p_count}, table={table_cell_para_count})"
            )
            return numbered_texts

        except ImportError:
            raise FileParseError("请安装python-docx库: pip install python-docx")
        except Exception as e:
            raise FileParseError(f"解析Word文档失败: {str(e)}")



class PowerPointParser(FileParser):
    """PowerPoint文档解析器"""
    
    def parse(self) -> Dict[str, Dict[str, Any]]:
        """解析PowerPoint文档"""
        try:
            from pptx import Presentation
            from zipfile import ZipFile
            from lxml import etree
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # 必须引入这个枚举
            
            prs = Presentation(self.file_path)
            numbered_texts = {}
            index = 1

            # --- 新增：递归遍历形状的辅助函数 ---
            def _iter_shapes(shapes):
                """递归生成器，能够进入Group和Table内部"""
                for shape in shapes:
                    # 1. 如果是组合形状 (Group)，递归遍历子形状
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from _iter_shapes(shape.shapes)
                    
                    # 2. 如果是表格 (Table)，遍历单元格
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                yield cell  # 单元格对象也有 text_frame 属性，用法同 shape
                                
                    # 3. 普通形状，直接返回
                    else:
                        yield shape
            # ---------------------------------------

            # 遍历所有幻灯片
            for slide_idx, slide in enumerate(prs.slides):
                # 使用递归函数遍历所有形状（包括组合内部和表格内部）
                for shape_idx, shape in enumerate(_iter_shapes(slide.shapes)):
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # 按段落分割
                            paragraphs = text.split('\n')
                            for para_idx, para_text in enumerate(paragraphs):
                                para_text = para_text.strip()
                                if para_text:
                                    number = f"S{slide_idx+1:02d}-P{index:04d}"
                                    # 注意：这里的 shape_idx 变成了扁平化后的索引，不仅是顶层索引
                                    # 但只要 Generator 使用相同的遍历顺序，就能对应上
                                    numbered_texts[number] = {
                                        'key': f'slide_{slide_idx}_flatshape_{shape_idx}_para_{para_idx}',
                                        'text': para_text,
                                        'type': 'slide_text',
                                        'slide_index': slide_idx,
                                        'paragraph_index': para_idx,
                                        'original_paragraphs': paragraphs
                                    }
                                    index += 1

            # OOXML：SmartArt（diagrams）、图表、备注页；与形状 API 并行，按部件路径+序号写回
            try:
                from core.ppt_ooxml import (
                    ppt_ooxml_list_writable_text_nodes,
                    ppt_ooxml_ordered_part_paths,
                )

                with ZipFile(self.file_path, "r") as zf:
                    for part_path in ppt_ooxml_ordered_part_paths(zf):
                        if part_path not in zf.namelist():
                            continue
                        xml_bytes = zf.read(part_path)
                        root = etree.fromstring(xml_bytes)
                        for t_idx, el in enumerate(
                            ppt_ooxml_list_writable_text_nodes(root, part_path)
                        ):
                            t = (el.text or "").strip()
                            number = f"S00-O{index:04d}"
                            numbered_texts[number] = {
                                "key": f"ppt_ooxml_{part_path}_{t_idx}",
                                "text": t,
                                "type": "ppt_ooxml_text",
                                "part_path": part_path,
                                "t_index": t_idx,
                                "slide_index": None,
                            }
                            index += 1
            except Exception as e:
                logger.warning(
                    "PPT OOXML 文本提取失败（若无 SmartArt/图表/备注可忽略）: %s", e
                )
            
            logger.info(f"PowerPoint文档解析完成，共提取 {len(numbered_texts)} 个文本段落")
            return numbered_texts
            
        except ImportError:
            raise FileParseError("请安装python-pptx库: pip install python-pptx")
        except Exception as e:
            raise FileParseError(f"解析PowerPoint文档失败: {str(e)}")


class ExcelParser(FileParser):
    """Excel文档解析器"""
    
    def parse(self) -> Dict[str, Dict[str, Any]]:
        """解析Excel文档"""
        try:
            from openpyxl import load_workbook
            # data_only=True 读取公式结果，keep_vba=True 保留宏
            # 注意: 如果是 xlsm，建议开启 keep_vba=True
            is_macro = self.file_path.lower().endswith('.xlsm')
            wb = load_workbook(self.file_path, data_only=True, keep_vba=is_macro)
            numbered_texts = {}
            index = 1
            
            # 遍历所有工作表
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # 遍历所有有值的单元格
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            text = str(cell.value).strip()
                            if text:
                                number = f"E{index:04d}"
                                numbered_texts[number] = {
                                    'key': f'sheet_{sheet_name}_cell_{cell.coordinate}',
                                    'text': text,
                                    'type': 'cell',
                                    'sheet_name': sheet_name,
                                    'cell_coordinate': cell.coordinate
                                }
                                index += 1
            
            logger.info(f"Excel文档解析完成，共提取 {len(numbered_texts)} 个文本单元格")
            return numbered_texts
            
        except ImportError:
            raise FileParseError("请安装openpyxl库: pip install openpyxl")
        except Exception as e:
            raise FileParseError(f"解析Excel文档失败: {str(e)}")


class PDFParser(FileParser):
    """PDF文档解析器"""
    
    def parse(self) -> Dict[str, Dict[str, Any]]:
        """解析PDF文档"""
        try:
            import pdfplumber
            
            numbered_texts = {}
            index = 1
            
            with pdfplumber.open(self.file_path) as pdf:
                # 遍历所有页面
                for page_idx, page in enumerate(pdf.pages):
                    try:
                        words = page.extract_words(
                            extra_attrs=["size", "fontname", "non_stroking_color"]
                        )
                    except Exception:
                        # 部分 PDF / pdfplumber 版本不支持颜色属性，回退到原行为
                        words = page.extract_words(extra_attrs=["size", "fontname"])
                    if not words:
                        continue

                    # 估计行/段聚合阈值（基于词的高度）
                    word_heights = [float(w.get("height", 0) or 0) for w in words]
                    word_heights = [h for h in word_heights if h > 0]
                    avg_word_h = (sum(word_heights) / len(word_heights)) if word_heights else 10.0
                    line_tol = max(1.5, avg_word_h * 0.25)

                    # 1) words -> lines
                    words_sorted = sorted(words, key=lambda w: (float(w.get("top", 0) or 0), float(w.get("x0", 0) or 0)))
                    lines = []
                    current_line = []
                    current_center = None

                    for w in words_sorted:
                        top = float(w.get("top", 0) or 0)
                        bottom = float(w.get("bottom", 0) or 0)
                        center = (top + bottom) / 2
                        if current_center is None:
                            current_line = [w]
                            current_center = center
                            continue

                        if abs(center - current_center) <= line_tol:
                            current_line.append(w)
                            # 更新行中心（让聚合更稳）
                            current_center = (
                                current_center * (len(current_line) - 1) + center
                            ) / len(current_line)
                        else:
                            lines.append(current_line)
                            current_line = [w]
                            current_center = center

                    if current_line:
                        lines.append(current_line)

                    # 2) lines -> paragraphs（根据段间垂直间隙）
                    line_objs = []
                    for line_words in lines:
                        line_words = sorted(line_words, key=lambda w: float(w.get("x0", 0) or 0))
                        if not line_words:
                            continue

                        # 同一行可能横跨多列：按较大的 x 间隙拆成多个“行块”
                        # 这样可避免多列内容被拼成一个大框后都写回到第一列。
                        split_gap = max(18.0, avg_word_h * 1.2)
                        chunks = []
                        cur_chunk = [line_words[0]]
                        prev_x1 = float(line_words[0].get("x1", 0) or 0)
                        for w in line_words[1:]:
                            x0_cur = float(w.get("x0", 0) or 0)
                            if x0_cur - prev_x1 > split_gap:
                                chunks.append(cur_chunk)
                                cur_chunk = [w]
                            else:
                                cur_chunk.append(w)
                            prev_x1 = float(w.get("x1", 0) or 0)
                        if cur_chunk:
                            chunks.append(cur_chunk)

                        for chunk_words in chunks:
                            x0 = min(float(w.get("x0", 0) or 0) for w in chunk_words)
                            x1 = max(float(w.get("x1", 0) or 0) for w in chunk_words)
                            y0 = min(float(w.get("top", 0) or 0) for w in chunk_words)
                            y1 = max(float(w.get("bottom", 0) or 0) for w in chunk_words)

                            text = " ".join(
                                str(w.get("text", "")).strip()
                                for w in chunk_words
                                if str(w.get("text", "")).strip()
                            )

                            sizes = [float(w.get("size", 0) or 0) for w in chunk_words]
                            sizes = [s for s in sizes if s > 0]
                            avg_size = (sum(sizes) / len(sizes)) if sizes else 10.0

                            # 尝试保留行级原文颜色（用于写回译文）
                            colors = [
                                w.get("non_stroking_color")
                                for w in chunk_words
                                if "non_stroking_color" in w
                            ]
                            line_color = None
                            if colors:
                                # 以“最后一个非空颜色”为准（常见场景下同块颜色一致）
                                for c in reversed(colors):
                                    if c is not None:
                                        line_color = c
                                        break

                            if text.strip():
                                line_objs.append(
                                    {
                                        "text": text.strip(),
                                        "rect": [x0, y0, x1, y1],
                                        "font_size": avg_size,
                                        "text_color": line_color,
                                        "top": y0,
                                        "bottom": y1,
                                    }
                                )

                    line_objs.sort(key=lambda lo: (lo["top"], lo["rect"][0]))

                    # 3) 直接按“行级”输出，避免段落级大 bbox 导致大面积白底遮挡
                    for line_idx, lo in enumerate(line_objs):
                        line_text = str(lo.get("text", "")).strip()
                        if not line_text:
                            continue

                        x0, y0, x1, y1 = [float(v) for v in lo["rect"]]
                        line_font_size = float(lo.get("font_size", 10.0) or 10.0)

                        number = f"Page{page_idx+1:02d}-P{index:04d}"
                        numbered_texts[number] = {
                            "key": f"page_{page_idx}_line_{line_idx}",
                            "text": line_text,
                            "type": "line",
                            "page_index": page_idx,
                            "paragraph_index": line_idx,
                            # (x0, y0, x1, y1) in top-left coordinate system
                            "rect": [x0, y0, x1, y1],
                            "font_size": line_font_size,
                            "text_color": lo.get("text_color"),
                        }
                        index += 1
            
            logger.info(f"PDF文档解析完成，共提取 {len(numbered_texts)} 个文本段落")
            return numbered_texts
            
        except ImportError:
            raise FileParseError("请安装pdfplumber库: pip install pdfplumber")
        except Exception as e:
            raise FileParseError(f"解析PDF文档失败: {str(e)}")


# === 修改 2: 完善 create_parser 映射表 ===
def create_parser(file_path: str) -> FileParser:
    """根据文件类型创建相应的解析器"""
    file_ext = Path(file_path).suffix.lower()
    
    parser_map = {
        # Word
        '.docx': WordParser,
        '.docm': WordParser,  # 新增
        
        # PPT
        '.pptx': PowerPointParser,
        '.pptm': PowerPointParser, # 新增
        
        # Excel
        '.xlsx': ExcelParser,
        '.xlsm': ExcelParser, # 新增
        # 注意：这里删除了 .xls，因为 parser 无法直接处理，要在 main 中转换
        
        # PDF
        '.pdf': PDFParser,
    }
    
    parser_class = parser_map.get(file_ext)
    
    # 针对转换后的临时文件处理 (比如 .doc -> _temp_converted.docx)
    if parser_class is None:
        raise FileParseError(f"不支持的文件类型: {file_ext} (请确保老格式文件已转换)")
    
    return parser_class(file_path)
